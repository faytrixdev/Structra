from __future__ import annotations

import asyncio
import io
import re
from typing import Optional

import PyPDF2
import docx
import markdown
import openpyxl
from bs4 import BeautifulSoup
from pptx import Presentation


PageCount = Optional[int]


async def extract_text(
    file_path: str, file_type: str, content: bytes
) -> tuple[str, PageCount]:
    if "pdf" in file_type:
        return await asyncio.to_thread(extract_pdf, content)
    elif "wordprocessingml" in file_type or "docx" in file_type:
        return await asyncio.to_thread(extract_docx, content)
    elif "spreadsheetml" in file_type or "xlsx" in file_type:
        return await asyncio.to_thread(extract_xlsx, content)
    elif "presentationml" in file_type or "pptx" in file_type:
        return await asyncio.to_thread(extract_pptx, content)
    elif "html" in file_type:
        return await asyncio.to_thread(extract_html, content)
    elif "xml" in file_type:
        return await asyncio.to_thread(extract_xml, content)
    elif "markdown" in file_type:
        return await asyncio.to_thread(extract_markdown, content)
    else:
        return content.decode("utf-8"), None


def extract_pdf(content: bytes) -> tuple[str, PageCount]:
    reader = PyPDF2.PdfReader(io.BytesIO(content))
    text: str = "\n".join(page.extract_text() for page in reader.pages)
    return text, len(reader.pages)


def extract_docx(content: bytes) -> tuple[str, PageCount]:
    doc = docx.Document(io.BytesIO(content))
    text: str = "\n".join(p.text for p in doc.paragraphs)
    return text, None


def extract_xlsx(content: bytes) -> tuple[str, PageCount]:
    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True)
    lines: list[str] = []
    for sheet in wb.worksheets:
        lines.append(f"\n## Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            lines.append(" | ".join(str(c) if c is not None else "" for c in row))
    return "\n".join(lines), None


def extract_pptx(content: bytes) -> tuple[str, PageCount]:
    prs = Presentation(io.BytesIO(content))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides):
        lines.append(f"\n## Slide {i + 1}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text)
    return "\n".join(lines), len(prs.slides)


def extract_html(content: bytes) -> tuple[str, PageCount]:
    text: str = BeautifulSoup(content, "html.parser").get_text(
        separator="\n", strip=True
    )
    return text, None


def extract_xml(content: bytes) -> tuple[str, PageCount]:
    text: str = BeautifulSoup(content, "xml").get_text(separator="\n", strip=True)
    return text, None


def extract_markdown(content: bytes) -> tuple[str, PageCount]:
    html: str = markdown.markdown(content.decode("utf-8"))
    text: str = BeautifulSoup(html, "html.parser").get_text(
        separator="\n", strip=True
    )
    return text, None


_SECTION_MARKDOWN_RE = re.compile(r"(?m)^## (.+)$")
_SECTION_PPTX_RE = re.compile(r"(?m)^## Slide \d+$")
_SECTION_XLSX_RE = re.compile(r"(?m)^## Sheet: .+$")
_SECTION_HTML_H1_RE = re.compile(r"(?s)(<h1\b[^>]*>.*?</h1>)")
_SECTION_HTML_H2_RE = re.compile(r"(?s)(<h2\b[^>]*>.*?</h2>)")
_NEUTRAL_HTML_TAG_RE = re.compile(r"<[^>]+>")


def _split_with_markers(
    text: str, marker_re: re.Pattern[str], extract_title
) -> list[tuple[str, str]]:
    sections: list[tuple[str, str]] = []
    matches: list[tuple[int, int, str, str]] = []

    for match in marker_re.finditer(text):
        title = extract_title(match.group(0))
        matches.append((match.start(), match.end(), match.group(0), title))

    if not matches:
        return [(text.strip(), "")] if text.strip() else []

    pre = text[: matches[0][0]].strip()
    if pre:
        sections.append(("", pre))

    for index, (start, end, raw, title) in enumerate(matches):
        body_start = end
        body_end = matches[index + 1][0] if index + 1 < len(matches) else len(text)
        body = text[body_start:body_end].strip()
        sections.append((title, body))

    return [s for s in sections if s[0] or s[1]]


def _markdown_title(match: str) -> str:
    return match[3:].strip()


def _pptx_title(match: str) -> str:
    return match.strip()


def _xlsx_title(match: str) -> str:
    return match[len("## Sheet:"):].strip()


def _sliding_windows(text: str, window_size: int = 1500, overlap: int = 200) -> list[tuple[str, str]]:
    stripped = text.strip()
    if not stripped:
        return []

    step = window_size - overlap
    if step <= 0:
        return [(f"Window 1", stripped)]

    sections: list[tuple[str, str]] = []
    start = 0
    total = len(stripped)
    index = 1
    while start < total:
        end = min(start + window_size, total)
        body = stripped[start:end].strip()
        if body:
            sections.append((f"Window {index}", body))
            index += 1
        if end == total:
            break
        start += step

    return sections


def _split_html(text: str) -> list[tuple[str, str]]:
    headings_re: re.Pattern[str]
    primary = _SECTION_HTML_H1_RE.findall(text)
    if primary:
        headings_re = _SECTION_HTML_H1_RE
    else:
        secondary = _SECTION_HTML_H2_RE.findall(text)
        if not secondary:
            return []
        headings_re = _SECTION_HTML_H2_RE

    sections: list[tuple[str, str]] = []
    cursor = 0
    while cursor < len(text):
        match = headings_re.search(text, cursor)
        if match is None:
            break
        next_match = headings_re.search(text, match.end())
        body_end = next_match.start() if next_match is not None else len(text)
        heading_html = match.group(0)
        body_raw = text[match.end():body_end]
        title = _NEUTRAL_HTML_TAG_RE.sub("", heading_html).strip()
        body = _NEUTRAL_HTML_TAG_RE.sub(" ", body_raw).strip()
        sections.append((title, body))
        cursor = body_end

    return [s for s in sections if s[0] or s[1]]


def detect_sections(text: str, file_type: str) -> list[tuple[str, str]]:
    normalized = (file_type or "").lower()
    sections: list[tuple[str, str]]

    if "markdown" in normalized:
        sections = _split_with_markers(text, _SECTION_MARKDOWN_RE, _markdown_title)
    elif "presentationml" in normalized or "pptx" in normalized:
        sections = _split_with_markers(text, _SECTION_PPTX_RE, _pptx_title)
    elif "spreadsheetml" in normalized or "xlsx" in normalized:
        sections = _split_with_markers(text, _SECTION_XLSX_RE, _xlsx_title)
    elif "html" in normalized:
        sections = _split_html(text)
        if not sections:
            sections = _sliding_windows(text)
    else:
        blank_split = re.split(r"\n{3,}", text)
        cleaned = [part.strip() for part in blank_split if part.strip()]
        if len(cleaned) > 1:
            sections = [("", part) for part in cleaned]
        else:
            sections = _sliding_windows(text)

    return sections
